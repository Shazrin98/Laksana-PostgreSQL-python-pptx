import psycopg2
from pptx import Presentation
from pptx.util import Inches

# Establish a connection to the PostgreSQL database
conn = psycopg2.connect(
    host="localhost",
    database="laporan_bulanan",
    user="postgres",
    password="shafiqnazrin"
)

# Create a cursor object to execute SQL commands
cur = conn.cursor()

# Fetch the data from the database
cur.execute("SELECT Nombor, Inisiatif, Agensi, PerbelanjaanRMJuta FROM laporan_bulanan LIMIT 10")
rows = cur.fetchall()

# Load the PowerPoint template
prs = Presentation('Laporan_Bulanan.pptx')

# Get the slide layout we want to use
slide_layout = prs.slide_layouts[0]

# Add a new slide using the layout
slide = prs.slides.add_slide(slide_layout)

# Add a table to the slide
table = slide.shapes.add_table(rows=6, cols=4, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)).table

# Add the table headers
table.cell(0, 0).text = 'Nombor'
table.cell(0, 1).text = 'Inisiatif'
table.cell(0, 2).text = 'Agensi'
table.cell(0, 3).text = 'Perbelanjaan (RM Juta)'

# Add the data rows to the table
for i, row in enumerate(rows):
    row_num = i + 1
    if row_num >= 6:
        # If there are more than 5 rows, add a new slide
        slide = prs.slides.add_slide(slide_layout)
        table = slide.shapes.add_table(rows=6, cols=4, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)).table
        table.cell(0, 0).text = 'Nombor'
        table.cell(0, 1).text = 'Inisiatif'
        table.cell(0, 2).text = 'Agensi'
        table.cell(0, 3).text = 'Perbelanjaan (RM Juta)'
        row_num = 1

    # Add the data to the table
    table.cell(row_num, 0).text = str(row[0])
    table.cell(row_num, 1).text = row[1]
    table.cell(row_num, 2).text = row[2]
    table.cell(row_num, 3).text = str(row[3])

# Save the PowerPoint presentation
prs.save('Laporan_Bulanan.pptx')

# Commit the changes to the database
conn.commit()

# Close the cursor and database connection
cur.close()
conn.close()




